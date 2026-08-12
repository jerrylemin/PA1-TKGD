from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

SOURCE_DIR = Path(__file__).resolve().parent
import os

RENDER_DIR = SOURCE_DIR.parent / "qa" / os.environ.get("RENDER_SUBDIR", "renders")
OUTPUT = SOURCE_DIR.parent / os.environ.get("OUTPUT_NAME", "Group10-PA3-PaperPrototypes.pptx")

slides = sorted(RENDER_DIR.glob("slide-*.png"))
if len(slides) != 20:
    raise SystemExit(f"Expected 20 rendered slides, found {len(slides)}")

for slide in slides:
    with Image.open(slide) as image:
        if image.size != (1920, 1080):
            raise SystemExit(f"Unexpected render size for {slide.name}: {image.size}")

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for slide_path in slides:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(slide_path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

prs.save(OUTPUT)
print(f"created {OUTPUT} with {len(prs.slides)} slides")
